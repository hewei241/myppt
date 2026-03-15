"""
生成 Diablo 4 风格 PPT 模板。
支持从 JSON 读取页面数据：首页为封面（cover.png），其余页为内容页（background.png）。
标题也可通过参数或命令行传入（仅生成封面时生效）。
"""

import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor

ASSETS_DIR = Path(__file__).parent / "assets" / "diablo4"
COVER_IMAGE = ASSETS_DIR / "cover.png"
BACKGROUND_IMAGE = ASSETS_DIR / "background.png"
DATA_JSON = ASSETS_DIR / "diablo4_class_heat_slides.json"
CLASS_IMAGE_DIR = ASSETS_DIR  # 职业图 class_职业名.png
OUTPUT_PATH = Path(__file__).parent / "diablo4_template.pptx"
TITLE_FONT_SIZE = Pt(60)
CONTENT_TITLE_FONT_SIZE = Pt(44)
CONTENT_BODY_FONT_SIZE = Pt(32)
DEFAULT_COVER_TITLE = "暗黑破坏神 IV"  # 未传参时的默认标题
COVER_TITLE_MAX_CHARS_PER_LINE = 4  # 每行最多显示字数


def wrap_title(text: str, max_chars: int = COVER_TITLE_MAX_CHARS_PER_LINE) -> str:
    """将标题按每行最多 max_chars 个字换行。"""
    return "\n".join(text[i : i + max_chars] for i in range(0, len(text), max_chars))


def get_class_image_path(class_name: str) -> Path | None:
    """根据职业名返回职业图路径（class_职业名.png 或 .webp），不存在则返回 None。"""
    if not class_name or not class_name.strip():
        return None
    name = class_name.strip()
    for ext in (".png", ".webp"):
        p = CLASS_IMAGE_DIR / f"class_{name}{ext}"
        if p.exists():
            return p
    return None


def title_to_class_name(title: str) -> str:
    """从幻灯片标题解析职业名，如 'No.7 德鲁伊' -> '德鲁伊'，'野蛮人' -> '野蛮人'。"""
    if not title:
        return ""
    t = title.strip()
    if t.startswith("No.") and " " in t:
        return t.split(" ", 1)[1].strip()
    return t


def gen_cover_slide(prs: Presentation, cover_path: Path, title: str) -> None:
    """添加封面页：背景图 + 居中标题。"""
    blank_layout = prs.slide_layouts[6]  # 6 = blank
    slide = prs.slides.add_slide(blank_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 1. 背景图（先添加，自然在底层）
    slide.shapes.add_picture(
        str(cover_path),
        Emu(0),
        Emu(0),
        width=slide_width,
        height=slide_height,
    )

    # 2. 居中标题（文本框高度覆盖整页）
    box_width = int(slide_width * 0.8)
    box_left = int((slide_width - box_width) / 2)
    box_top = 0
    box_height = slide_height

    tx = slide.shapes.add_textbox(box_left, box_top, box_width, box_height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # 标题垂直居中
    p = tf.paragraphs[0]
    p.text = wrap_title(title)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "微软雅黑"
    p.font.size = TITLE_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 140, 0)  # 橙色


def gen_content_slide(
    prs: Presentation,
    bg_path: Path,
    title: str,
    text: str,
    class_name: str | None = None,
) -> None:
    """添加内容页：background.png 为底图，左侧标题+正文，右侧可选职业图。"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 0. 幻灯片背景设为黑色，避免右侧或边缘露出白色
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    # 1. 背景图（铺满整页，使用整数 EMU 避免缝隙）
    slide.shapes.add_picture(
        str(bg_path),
        Emu(0),
        Emu(0),
        width=slide_width,
        height=slide_height,
    )

    # 2. 标题（紧贴左上角，44 号字）
    margin = int(slide_width * 0.02)
    title_top = int(slide_height * 0.02)
    title_width = int(slide_width * 0.5) - margin
    title_height = Emu(600000)

    tx_title = slide.shapes.add_textbox(margin, title_top, title_width, title_height)
    tf_title = tx_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = "微软雅黑"
    p_title.font.size = CONTENT_TITLE_FONT_SIZE
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 140, 0)

    # 3. 正文（标题下方，仅占左半边，支持 \n 换行）
    body_top = title_top + int(slide_height * 0.10)
    body_width = int(slide_width * 0.5) - margin
    body_height = int(slide_height * 0.7)

    tx_body = slide.shapes.add_textbox(margin, body_top, body_width, body_height)
    tf_body = tx_body.text_frame
    tf_body.word_wrap = True
    tf_body.auto_size = None
    for i, line in enumerate(text.split("\n")):
        p = tf_body.paragraphs[0] if i == 0 else tf_body.add_paragraph()
        p.text = line
        p.font.name = "微软雅黑"
        p.font.size = CONTENT_BODY_FONT_SIZE
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(12)

    # 4. 右侧职业图（若存在 class_职业名.png / .webp）
    class_img = get_class_image_path(class_name or title_to_class_name(title))
    if class_img is not None:
        img_width = int(slide_width * 0.42)
        img_left = int(slide_width * 0.54)
        img_top = int(slide_height * 0.15)
        img_height = int(slide_height * 0.7)
        slide.shapes.add_picture(
            str(class_img),
            Emu(img_left),
            Emu(img_top),
            width=img_width,
            height=img_height,
        )


def load_slides_data(json_path: Path) -> list[dict] | None:
    """从 JSON 加载页面数据列表。"""
    if not json_path.exists():
        return None
    try:
        with open(json_path, "r", encoding="utf-8") as f:
            return json.load(f)
    except (json.JSONDecodeError, OSError):
        return None


def main(title: str | None = None, use_json: bool = True):
    if not COVER_IMAGE.exists():
        print(f"请将封面图放在: {COVER_IMAGE}")
        print("  或修改脚本顶部 COVER_IMAGE 路径。")
        exit(1)

    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 16:9 约 13.333 英寸
    prs.slide_height = Emu(6858000)   # 约 7.5 英寸

    slides_data = load_slides_data(DATA_JSON) if use_json else None

    if slides_data:
        for item in slides_data:
            if item.get("type") == "cover":
                gen_cover_slide(prs, COVER_IMAGE, item.get("title", DEFAULT_COVER_TITLE))
            elif item.get("type") == "slide":
                if not BACKGROUND_IMAGE.exists():
                    print(f"请将内容页背景图放在: {BACKGROUND_IMAGE}")
                    exit(1)
                slide_title = item.get("title", "")
                gen_content_slide(
                    prs,
                    BACKGROUND_IMAGE,
                    slide_title,
                    item.get("text", ""),
                    class_name=title_to_class_name(slide_title),
                )
        print(f"已根据 {DATA_JSON} 生成 {len(slides_data)} 页 → {OUTPUT_PATH}")
    else:
        cover_title = title if title is not None else DEFAULT_COVER_TITLE
        gen_cover_slide(prs, COVER_IMAGE, cover_title)
        print(f"已生成: {OUTPUT_PATH}")

    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    # 支持命令行传参：python gen_diablo4_ppt.py "我的标题"（仅无 JSON 时生效）
    title_arg = sys.argv[1] if len(sys.argv) > 1 else None
    main(title=title_arg)
